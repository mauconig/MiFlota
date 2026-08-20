from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "presentaciones"

COLORS = {
    "ink": "16150F",
    "paper": "F4F0E8",
    "white": "FFFDF8",
    "amber": "E8A13A",
    "amber_light": "F7DCA8",
    "gold": "8A5D16",
    "green": "256B4D",
    "green_light": "D9E8DC",
    "coral": "914531",
    "coral_light": "F3D9D0",
    "muted": "746F65",
    "line": "D8D0C2",
}

FONT = "Plus Jakarta Sans"
WIDTH = 13.333
HEIGHT = 7.5


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def fill(shape, color: str, transparency: int = 0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(COLORS[color] if color in COLORS else color)
    shape.fill.transparency = transparency


def no_line(shape):
    shape.line.fill.background()


def line(shape, color: str, width: float = 1):
    shape.line.color.rgb = rgb(COLORS[color] if color in COLORS else color)
    shape.line.width = Pt(width)


def add_box(slide, x, y, w, h, color="white", radius=True, border=None, border_width=1):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    fill(shape, color)
    if border:
        line(shape, border, border_width)
    else:
        no_line(shape)
    return shape


def add_line(slide, x1, y1, x2, y2, color="line", width=1, dashed=False):
    shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line(shape, color, width)
    if dashed:
        shape.line.dash_style = 1
    return shape


def add_text(slide, text, x, y, w, h, size=16, color="ink", bold=False, align=PP_ALIGN.LEFT,
             valign=MSO_ANCHOR.TOP, font=FONT, margin=0.04, italic=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(COLORS[color] if color in COLORS else color)
    return box


def add_rich_lines(slide, lines, x, y, w, h, size=15, color="ink", bullet_color="amber", gap=0.38):
    """Draw a compact list where each line is a bullet and body text."""
    for idx, item in enumerate(lines):
        yy = y + idx * gap
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(yy + 0.09), Inches(0.1), Inches(0.1))
        fill(dot, bullet_color)
        no_line(dot)
        add_text(slide, item, x + 0.18, yy, w - 0.18, 0.28, size=size, color=color)


def add_mark(slide, x, y, size=0.42, dark=False):
    shape = add_box(slide, x, y, size, size, "ink" if dark else "amber", radius=True)
    add_text(slide, "M", x, y + 0.005, size, size - 0.01, size=16 if size < 0.6 else 24,
             color="white" if dark else "ink", bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    return shape


def add_label(slide, text, x, y, w=2.2, color="gold"):
    add_text(slide, text.upper(), x, y, w, 0.18, size=8.5, color=color, bold=True, margin=0,
             font="Aptos Display")


def add_title(slide, kicker, title, subtitle=None, page=None, dark=False):
    add_mark(slide, 0.55, 0.42, 0.38, dark=dark)
    add_label(slide, kicker, 1.05, 0.48, 3.8, color="amber" if dark else "gold")
    add_text(slide, title, 0.55, 0.95, 8.8, 0.72, size=29, color="white" if dark else "ink", bold=True, margin=0)
    if subtitle:
        add_text(slide, subtitle, 0.58, 1.72, 8.8, 0.44, size=12.5, color="amber_light" if dark else "muted", margin=0)
    if page:
        add_text(slide, f"{page:02d}", 12.0, 0.48, 0.65, 0.22, size=9, color="amber_light" if dark else "muted",
                 align=PP_ALIGN.RIGHT, margin=0)


def add_footer(slide, product, page):
    add_line(slide, 0.55, 7.08, 12.78, 7.08, color="line", width=0.6)
    add_text(slide, f"MiFlota  ·  {product}", 0.58, 7.16, 4.4, 0.18, size=8.5, color="muted", margin=0)
    add_text(slide, f"{page:02d}", 12.1, 7.16, 0.6, 0.18, size=8.5, color="muted", align=PP_ALIGN.RIGHT, margin=0)


def base_slide(prs, dark=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb(COLORS["ink"] if dark else COLORS["paper"])
    return slide


def add_pill(slide, text, x, y, w, color="amber", text_color="ink", h=0.34, size=9.5):
    add_box(slide, x, y, w, h, color, radius=True)
    add_text(slide, text, x + 0.02, y + 0.005, w - 0.04, h - 0.01, size=size, color=text_color,
             bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0)


def add_metric(slide, x, y, w, value, label, accent="amber", dark=False):
    card = add_box(slide, x, y, w, 1.14, "ink" if dark else "white", radius=True,
                   border=None if dark else "line")
    add_box(slide, x, y, 0.08, 1.14, accent, radius=False)
    add_text(slide, value, x + 0.24, y + 0.18, w - 0.35, 0.42, size=24, color="white" if dark else "ink", bold=True, margin=0)
    add_text(slide, label, x + 0.24, y + 0.67, w - 0.35, 0.24, size=9.5, color="amber_light" if dark else "muted", margin=0)
    return card


def add_browser(slide, x, y, w, h, accent="amber", title="Resumen"):
    # Outer browser frame.
    add_box(slide, x, y, w, h, "ink", radius=True)
    add_box(slide, x + 0.12, y + 0.34, w - 0.24, h - 0.46, "white", radius=False)
    for idx, c in enumerate(["coral", "amber", "green"]):
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.18 + idx * 0.16), Inches(y + 0.12), Inches(0.09), Inches(0.09))
        fill(dot, c)
        no_line(dot)
    add_text(slide, "miflota / panel", x + 0.8, y + 0.075, w - 1.4, 0.18, size=8.5, color="amber_light", margin=0)
    # Sidebar and page heading.
    add_box(slide, x + 0.12, y + 0.34, 1.26, h - 0.46, "ink", radius=False)
    add_mark(slide, x + 0.33, y + 0.58, 0.29, dark=False)
    for i, label in enumerate(["Resumen", "Flota", "Choferes", "Alertas", "Reportes"]):
        yy = y + 1.12 + i * 0.43
        if i == 0:
            add_box(slide, x + 0.23, yy - 0.03, 1.02, 0.28, accent, radius=True)
        add_text(slide, label, x + 0.36, yy + 0.01, 0.8, 0.16, size=7.2, color="ink" if i == 0 else "paper", margin=0)
    add_text(slide, title, x + 1.65, y + 0.58, 2.7, 0.28, size=14, color="ink", bold=True, margin=0)
    add_text(slide, "Actualizado hace 2 min", x + 1.65, y + 0.88, 2.6, 0.16, size=7.5, color="muted", margin=0)
    # Dashboard cards.
    card_y = y + 1.2
    add_metric(slide, x + 1.65, card_y, 1.44, "$ 18.4M", "neto del mes", accent=accent)
    add_metric(slide, x + 3.2, card_y, 1.44, "12", "vehículos", accent="green")
    add_metric(slide, x + 4.75, card_y, 1.12, "86%", "cobrado", accent="coral")
    # Chart area.
    add_box(slide, x + 1.65, y + 2.6, 2.8, 1.33, "paper", radius=True, border="line")
    add_text(slide, "Ingresos vs. gastos", x + 1.83, y + 2.78, 1.9, 0.16, size=8.5, color="ink", bold=True, margin=0)
    points = [(x + 1.88, y + 3.52), (x + 2.25, y + 3.18), (x + 2.65, y + 3.35), (x + 3.02, y + 2.98), (x + 3.41, y + 3.16), (x + 3.83, y + 2.76), (x + 4.25, y + 2.91)]
    for a, b in zip(points, points[1:]):
        add_line(slide, a[0], a[1], b[0], b[1], color=accent, width=2)
    for px, py in points:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(px - 0.04), Inches(py - 0.04), Inches(0.08), Inches(0.08))
        fill(dot, accent)
        no_line(dot)
    add_box(slide, x + 4.62, y + 2.6, 1.25, 1.33, "paper", radius=True, border="line")
    add_text(slide, "Alertas", x + 4.8, y + 2.78, 0.78, 0.16, size=8.5, color="ink", bold=True, margin=0)
    add_text(slide, "03", x + 4.78, y + 3.05, 0.8, 0.28, size=22, color="coral", bold=True, margin=0)
    add_text(slide, "requieren acción", x + 4.8, y + 3.42, 0.9, 0.16, size=7.3, color="muted", margin=0)


def add_phone(slide, x, y, w=2.55, h=5.0, accent="amber", driver=False):
    add_box(slide, x, y, w, h, "ink", radius=True)
    add_box(slide, x + 0.13, y + 0.17, w - 0.26, h - 0.31, "white", radius=True)
    add_box(slide, x + w / 2 - 0.28, y + 0.07, 0.56, 0.12, "ink", radius=True)
    # Header
    add_mark(slide, x + 0.29, y + 0.43, 0.3, dark=False)
    add_text(slide, "MiFlota", x + 0.67, y + 0.47, 1.1, 0.18, size=8.5, color="ink", bold=True, margin=0)
    if driver:
        add_pill(slide, "CHOFER", x + 1.48, y + 0.46, 0.62, color="green_light", text_color="green", h=0.24, size=7)
    else:
        add_pill(slide, "ADMIN", x + 1.48, y + 0.46, 0.62, color="amber_light", text_color="gold", h=0.24, size=7)
    # Content cards.
    add_text(slide, "Hola, Diego" if driver else "Buen día, Martín", x + 0.3, y + 0.92, w - 0.6, 0.26, size=13, color="ink", bold=True, margin=0)
    add_text(slide, "Tu resumen de hoy" if driver else "Tu flota en un vistazo", x + 0.31, y + 1.22, w - 0.6, 0.17, size=7.5, color="muted", margin=0)
    add_box(slide, x + 0.3, y + 1.62, w - 0.6, 0.9, accent, radius=True)
    add_text(slide, "$ 1.080.000" if driver else "$ 18.4M", x + 0.47, y + 1.78, w - 0.9, 0.3, size=18, color="ink", bold=True, margin=0)
    add_text(slide, "deuda pendiente" if driver else "neto del mes", x + 0.48, y + 2.14, w - 0.9, 0.18, size=7.5, color="ink", margin=0)
    if driver:
        rows = [("Próximo vencimiento", "07 ago 2026", "coral"), ("Pagos registrados", "14", "green"), ("Reportar falla", "→", "gold")]
    else:
        rows = [("Cobros", "$ 4.2M", "green"), ("Gastos", "$ 1.6M", "coral"), ("Salud de flota", "92%", "gold")]
    for i, (label, value, color) in enumerate(rows):
        yy = y + 2.88 + i * 0.52
        add_box(slide, x + 0.3, yy, w - 0.6, 0.38, "paper", radius=True, border="line")
        add_text(slide, label, x + 0.43, yy + 0.1, 1.3, 0.16, size=7.2, color="muted", margin=0)
        add_text(slide, value, x + 1.65, yy + 0.1, 0.55, 0.16, size=7.5, color=color, bold=True, align=PP_ALIGN.RIGHT, margin=0)
    add_box(slide, x + 0.3, y + h - 0.67, w - 0.6, 0.34, "ink", radius=True)
    add_text(slide, "Inicio     Flota     Más", x + 0.39, y + h - 0.58, w - 0.78, 0.15, size=7, color="paper", align=PP_ALIGN.CENTER, margin=0)


def add_arrow(slide, x1, y1, x2, y2, color="amber"):
    shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line(shape, color, 2)
    shape.line.end_arrowhead = True
    return shape


def add_feature_card(slide, x, y, w, h, number, title, body, accent="amber", dark=False):
    add_box(slide, x, y, w, h, "ink" if dark else "white", radius=True, border=None if dark else "line")
    add_pill(slide, number, x + 0.2, y + 0.2, 0.42, color=accent, text_color="ink", h=0.28, size=8)
    add_text(slide, title, x + 0.2, y + 0.65, w - 0.4, 0.28, size=13, color="white" if dark else "ink", bold=True, margin=0)
    add_text(slide, body, x + 0.2, y + 1.02, w - 0.4, h - 1.16, size=9.5, color="amber_light" if dark else "muted", margin=0)


def add_architecture(slide, product, scope, accent="amber"):
    # Three-tier architecture diagram.
    add_text(slide, "Una misma fuente de verdad para toda la operación", 0.58, 1.93, 6.8, 0.28, size=15, color="ink", bold=True, margin=0)
    nodes = [
        (0.72, 3.05, 2.35, 1.22, "Experiencia", scope, "white"),
        (4.02, 3.05, 2.35, 1.22, "API segura", "Fastify · sesiones · roles", "ink"),
        (7.32, 3.05, 2.35, 1.22, "Datos", "SQLite · WAL · auditoría", "white"),
    ]
    for x, y, w, h, title, body, color in nodes:
        add_box(slide, x, y, w, h, color, radius=True, border=None if color == "ink" else "line")
        add_text(slide, title, x + 0.2, y + 0.23, w - 0.4, 0.25, size=13, color="white" if color == "ink" else "ink", bold=True, margin=0)
        add_text(slide, body, x + 0.2, y + 0.61, w - 0.4, 0.38, size=9.2, color="amber_light" if color == "ink" else "muted", margin=0)
    add_arrow(slide, 3.15, 3.66, 4.0, 3.66, color=accent)
    add_arrow(slide, 6.45, 3.66, 7.3, 3.66, color=accent)
    add_text(slide, "HTTPS", 3.36, 3.29, 0.52, 0.18, size=7.5, color="gold", bold=True, align=PP_ALIGN.CENTER, margin=0)
    add_text(slide, "derivados", 6.65, 3.29, 0.72, 0.18, size=7.5, color="gold", bold=True, align=PP_ALIGN.CENTER, margin=0)
    # Data examples.
    add_pill(slide, "Cobros", 8.0, 4.65, 0.88, color="green_light", text_color="green", h=0.3, size=8)
    add_pill(slide, "Gastos", 8.98, 4.65, 0.86, color="coral_light", text_color="coral", h=0.3, size=8)
    add_pill(slide, "Alertas", 9.94, 4.65, 0.9, color="amber_light", text_color="gold", h=0.3, size=8)
    add_text(slide, "El mismo registro alimenta reportes, saldos, ranking y asistente IA.", 0.73, 5.15, 7.2, 0.28, size=11, color="muted", margin=0)


def deck_setup():
    prs = Presentation()
    prs.slide_width = Inches(WIDTH)
    prs.slide_height = Inches(HEIGHT)
    return prs


def web_deck():
    prs = deck_setup()
    # 1 Cover
    slide = base_slide(prs, dark=True)
    add_mark(slide, 0.7, 0.7, 0.66, dark=False)
    add_label(slide, "Presentación de producto", 1.55, 0.78, 4.0, color="amber_light")
    add_text(slide, "MiFlota Web", 0.7, 1.78, 7.1, 0.8, size=42, color="white", bold=True, margin=0)
    add_text(slide, "El centro de control para cada vehículo, cobro y decisión.", 0.74, 2.75, 6.8, 0.54, size=18, color="amber_light", margin=0)
    add_pill(slide, "Panel para dueños de flota", 0.74, 3.65, 2.35, color="amber", text_color="ink", h=0.4, size=10)
    add_browser(slide, 7.45, 1.05, 5.25, 4.83, accent="amber", title="Resumen")
    add_text(slide, "CONTROL · CLARIDAD · ACCIÓN", 0.74, 6.52, 4.2, 0.2, size=9, color="amber", bold=True, margin=0)
    add_footer(slide, "Web", 1)

    # 2 Problem/value
    slide = base_slide(prs)
    add_title(slide, "01 · La oportunidad", "De planillas dispersas a una sola vista", "MiFlota convierte la operación diaria en decisiones simples.", 2)
    add_feature_card(slide, 0.65, 2.45, 3.75, 2.15, "01", "Ver qué pasa", "Ingresos, gastos, deudas y salud de la flota en un dashboard vivo.", accent="amber")
    add_feature_card(slide, 4.78, 2.45, 3.75, 2.15, "02", "Entender por qué", "Filtros, reportes y movimientos permiten encontrar el origen de cada número.", accent="green")
    add_feature_card(slide, 8.91, 2.45, 3.75, 2.15, "03", "Actuar a tiempo", "Alertas y cobros registrados desde el mismo lugar, con trazabilidad.", accent="coral")
    add_text(slide, "Una flota rentable no se gestiona mirando hacia atrás: se gestiona con contexto.", 0.68, 5.35, 8.9, 0.42, size=19, color="ink", bold=True, margin=0)
    add_metric(slide, 10.1, 5.05, 2.3, "1", "fuente de verdad", accent="amber")
    add_footer(slide, "Web", 2)

    # 3 Workflow
    slide = base_slide(prs, dark=True)
    add_title(slide, "02 · Flujo de trabajo", "Una operación completa, de punta a punta", "Cada acción deja un registro útil para la siguiente decisión.", 3, dark=True)
    steps = [
        (0.75, "Registrar", "Cobro o gasto", "green"),
        (3.25, "Consolidar", "Movimiento + comprobante", "amber"),
        (5.75, "Analizar", "Saldos, alertas, ranking", "coral"),
        (8.25, "Decidir", "Priorizar la próxima acción", "green"),
    ]
    for i, (x, title, body, accent) in enumerate(steps):
        add_box(slide, x, 3.0, 2.0, 1.38, "white", radius=True)
        add_pill(slide, f"0{i+1}", x + 0.2, 3.22, 0.46, color=accent, text_color="ink", h=0.3, size=8)
        add_text(slide, title, x + 0.2, 3.66, 1.6, 0.22, size=13, color="ink", bold=True, margin=0)
        add_text(slide, body, x + 0.2, 4.0, 1.6, 0.2, size=8.5, color="muted", margin=0)
        if i < len(steps) - 1:
            add_arrow(slide, x + 2.08, 3.69, x + 2.42, 3.69, color="amber")
    add_text(slide, "De un movimiento individual a una lectura estratégica de toda la flota.", 0.78, 5.35, 8.8, 0.3, size=16, color="amber_light", margin=0)
    add_pill(slide, "menos fricción", 9.45, 5.22, 1.45, color="green_light", text_color="green", h=0.36, size=9)
    add_pill(slide, "más control", 11.05, 5.22, 1.3, color="amber_light", text_color="gold", h=0.36, size=9)
    add_footer(slide, "Web", 3)

    # 4 Screens
    slide = base_slide(prs)
    add_title(slide, "03 · Producto", "El panel que organiza el día del dueño", "Un mapa de navegación corto, visual y accionable.", 4)
    add_browser(slide, 0.68, 2.15, 7.0, 4.4, accent="amber", title="Resumen")
    add_box(slide, 8.3, 2.3, 4.25, 3.98, "ink", radius=True)
    add_label(slide, "Módulos clave", 8.65, 2.65, 2.5, color="amber")
    add_rich_lines(slide, [
        "Resumen: KPIs, rendimiento y alertas.",
        "Flota: tabla sortable y detalle por vehículo.",
        "Choferes: deuda, pagos y asignaciones.",
        "Reportes: movimientos, categorías y Excel.",
        "Cobros: cuotas, pagos y comprobantes.",
    ], 8.65, 3.15, 3.35, 2.25, size=10, color="white", bullet_color="amber", gap=0.44)
    add_pill(slide, "desktop-first", 8.65, 5.75, 1.35, color="amber", text_color="ink", h=0.33, size=8.5)
    add_pill(slide, "datos reales", 10.15, 5.75, 1.28, color="green_light", text_color="green", h=0.33, size=8.5)
    add_footer(slide, "Web", 4)

    # 5 Architecture
    slide = base_slide(prs)
    add_title(slide, "04 · Plataforma", "Una base sólida para crecer sin perder el control", "Seguridad y consistencia detrás de cada pantalla.", 5)
    add_architecture(slide, "Web", "Panel React + TypeScript", accent="amber")
    add_metric(slide, 10.15, 2.1, 2.2, "HTTPS", "dominio público", accent="green")
    add_metric(slide, 10.15, 3.48, 2.2, "RBAC", "sesiones por rol", accent="amber")
    add_metric(slide, 10.15, 4.86, 2.2, "WAL", "escrituras confiables", accent="coral")
    add_footer(slide, "Web", 5)

    # 6 Close
    slide = base_slide(prs, dark=True)
    add_title(slide, "05 · Demo", "La próxima decisión está a un clic", "MiFlota Web hace visible lo importante para que el equipo actúe.", 6, dark=True)
    add_box(slide, 0.75, 2.45, 6.0, 2.65, "white", radius=True)
    add_label(slide, "Recorrido sugerido", 1.12, 2.82, 2.5, color="gold")
    add_rich_lines(slide, ["Abrir Resumen y detectar una alerta.", "Entrar al vehículo o chofer involucrado.", "Registrar el cobro y exportar el reporte."], 1.12, 3.32, 5.1, 1.25, size=12, color="ink", bullet_color="amber", gap=0.47)
    add_browser(slide, 7.65, 1.86, 4.7, 4.25, accent="green", title="Cobros")
    add_text(slide, "De la señal al movimiento registrado.", 0.78, 5.82, 5.8, 0.28, size=16, color="amber_light", bold=True, margin=0)
    add_footer(slide, "Web", 6)
    return prs


def mobile_deck():
    prs = deck_setup()
    # 1 Cover
    slide = base_slide(prs, dark=True)
    add_mark(slide, 0.7, 0.7, 0.66, dark=False)
    add_label(slide, "Presentación de producto", 1.55, 0.78, 4.0, color="amber_light")
    add_text(slide, "MiFlota Admin", 0.7, 1.78, 7.2, 0.8, size=40, color="white", bold=True, margin=0)
    add_text(slide, "El control de la flota, en el bolsillo del dueño.", 0.74, 2.75, 6.8, 0.54, size=18, color="amber_light", margin=0)
    add_pill(slide, "App móvil para administradores", 0.74, 3.65, 2.7, color="amber", text_color="ink", h=0.4, size=10)
    add_phone(slide, 8.35, 1.03, w=3.0, h=5.95, accent="amber", driver=False)
    add_text(slide, "VISTA · REGISTRO · ASISTENTE", 0.74, 6.52, 4.5, 0.2, size=9, color="amber", bold=True, margin=0)
    add_footer(slide, "Admin mobile", 1)

    # 2 Value
    slide = base_slide(prs)
    add_title(slide, "01 · La oportunidad", "Administrar sin esperar a estar frente a una PC", "La información clave acompaña al dueño durante toda la jornada.", 2)
    add_phone(slide, 0.83, 1.95, w=2.65, h=5.25, accent="amber", driver=False)
    add_feature_card(slide, 4.2, 2.05, 3.8, 1.55, "01", "Resumen en segundos", "Neto, cobros, gastos y salud de flota en una pantalla.", accent="amber")
    add_feature_card(slide, 4.2, 3.83, 3.8, 1.55, "02", "Registrar donde ocurre", "Cobros y gastos con comprobante, sin volver a la oficina.", accent="green")
    add_feature_card(slide, 8.42, 2.05, 3.8, 1.55, "03", "Encontrar lo importante", "Búsqueda global, filtros y acceso directo al detalle.", accent="coral")
    add_feature_card(slide, 8.42, 3.83, 3.8, 1.55, "04", "Preguntar a MiFlota IA", "Respuestas con tarjetas, tablas y filtros tocables.", accent="amber")
    add_text(slide, "La app no reemplaza el criterio del dueño: le da contexto antes.", 4.23, 5.95, 7.8, 0.3, size=16, color="ink", bold=True, margin=0)
    add_footer(slide, "Admin mobile", 2)

    # 3 Workflow
    slide = base_slide(prs, dark=True)
    add_title(slide, "02 · Flujo diario", "Cinco gestos para mantener la operación al día", "Diseñada para consultar, registrar y volver a la ruta.", 3, dark=True)
    steps = [
        (0.72, "Abrir", "Dashboard", "amber"),
        (3.25, "Revisar", "alerta o saldo", "coral"),
        (5.78, "Registrar", "cobro / gasto", "green"),
        (8.31, "Filtrar", "detalle exacto", "amber"),
        (10.84, "Preguntar", "MiFlota IA", "green"),
    ]
    for i, (x, title, body, accent) in enumerate(steps):
        add_box(slide, x, 3.02, 1.78, 1.4, "white", radius=True)
        add_pill(slide, f"0{i+1}", x + 0.18, 3.23, 0.42, color=accent, text_color="ink", h=0.29, size=8)
        add_text(slide, title, x + 0.18, 3.68, 1.42, 0.22, size=12, color="ink", bold=True, margin=0)
        add_text(slide, body, x + 0.18, 4.01, 1.42, 0.2, size=8.1, color="muted", margin=0)
        if i < len(steps) - 1:
            add_arrow(slide, x + 1.86, 3.73, x + 2.45, 3.73, color="amber")
    add_text(slide, "La app mantiene la velocidad de la operación sin sacrificar trazabilidad.", 0.75, 5.35, 8.8, 0.3, size=16, color="amber_light", margin=0)
    add_pill(slide, "mobile-first", 9.65, 5.22, 1.45, color="amber_light", text_color="gold", h=0.36, size=9)
    add_footer(slide, "Admin mobile", 3)

    # 4 Screens
    slide = base_slide(prs)
    add_title(slide, "03 · Producto", "Una interfaz pensada para tocar y decidir", "Jerarquía clara, acciones directas y detalle cuando hace falta.", 4)
    add_phone(slide, 0.68, 1.85, w=2.48, h=5.05, accent="amber", driver=False)
    add_phone(slide, 3.65, 1.85, w=2.48, h=5.05, accent="green", driver=False)
    add_phone(slide, 6.62, 1.85, w=2.48, h=5.05, accent="coral", driver=False)
    add_box(slide, 9.72, 2.03, 2.58, 4.7, "ink", radius=True)
    add_label(slide, "Momentos clave", 10.05, 2.4, 1.8, color="amber")
    add_rich_lines(slide, ["Dashboard: entender el estado.", "Flota: entrar al detalle.", "Registrar: dejar evidencia.", "Ranking: comparar rendimiento.", "IA: preguntar con lenguaje natural."], 10.03, 2.93, 1.95, 2.5, size=9.5, color="white", bullet_color="amber", gap=0.48)
    add_footer(slide, "Admin mobile", 4)

    # 5 AI
    slide = base_slide(prs, dark=True)
    add_title(slide, "04 · MiFlota IA", "Preguntas simples, respuestas accionables", "La respuesta no termina en texto: abre el siguiente paso.", 5, dark=True)
    add_box(slide, 0.75, 2.15, 5.0, 3.45, "white", radius=True)
    add_label(slide, "Ejemplo", 1.1, 2.48, 1.2, color="gold")
    add_text(slide, "¿Cuánto debe Diego Rotela?", 1.1, 2.95, 3.9, 0.28, size=14, color="ink", bold=True, margin=0)
    add_box(slide, 1.1, 3.5, 4.3, 1.15, "paper", radius=True, border="line")
    add_text(slide, "Diego Rotela debe $ 1.080.000.", 1.32, 3.74, 3.8, 0.23, size=13, color="ink", bold=True, margin=0)
    add_text(slide, "La deuda más antigua es del 07 ago 2026.", 1.32, 4.14, 3.85, 0.2, size=9.2, color="muted", margin=0)
    add_pill(slide, "ver conductor", 1.1, 5.02, 1.28, color="amber", text_color="ink", h=0.32, size=8)
    add_pill(slide, "ver vehículo", 2.53, 5.02, 1.28, color="green_light", text_color="green", h=0.32, size=8)
    # Table / filters panel.
    add_box(slide, 6.35, 2.15, 5.95, 3.45, "ink", radius=True)
    add_label(slide, "Respuesta compuesta", 6.72, 2.48, 2.2, color="amber")
    add_text(slide, "La IA entiende la consulta y devuelve contexto.", 6.72, 2.91, 4.8, 0.22, size=13, color="white", bold=True, margin=0)
    add_box(slide, 6.72, 3.45, 4.95, 0.72, "white", radius=True)
    add_text(slide, "Chofer", 6.95, 3.65, 1.15, 0.17, size=8.3, color="muted", bold=True, margin=0)
    add_text(slide, "Deuda", 8.65, 3.65, 1.15, 0.17, size=8.3, color="muted", bold=True, margin=0)
    add_text(slide, "Antigüedad", 10.0, 3.65, 1.45, 0.17, size=8.3, color="muted", bold=True, margin=0)
    add_text(slide, "D. Rotela", 6.95, 3.94, 1.3, 0.17, size=9.5, color="ink", bold=True, margin=0)
    add_text(slide, "$ 1.08M", 8.65, 3.94, 1.3, 0.17, size=9.5, color="coral", bold=True, margin=0)
    add_text(slide, "07 ago", 10.0, 3.94, 1.2, 0.17, size=9.5, color="ink", margin=0)
    add_pill(slide, "deudores", 6.72, 4.65, 0.95, color="coral_light", text_color="coral", h=0.3, size=8)
    add_pill(slide, "agosto", 7.82, 4.65, 0.83, color="amber_light", text_color="gold", h=0.3, size=8)
    add_pill(slide, "prioridad alta", 8.8, 4.65, 1.2, color="green_light", text_color="green", h=0.3, size=8)
    add_footer(slide, "Admin mobile", 5)

    # 6 Architecture / close
    slide = base_slide(prs)
    add_title(slide, "05 · Plataforma", "Una app liviana, conectada al backend real", "La experiencia móvil comparte datos y reglas con el panel web.", 6)
    add_architecture(slide, "Admin mobile", "Expo · Android · iOS", accent="amber")
    add_metric(slide, 10.15, 2.1, 2.2, "1 API", "datos consistentes", accent="amber")
    add_metric(slide, 10.15, 3.48, 2.2, "IA", "respuesta + acción", accent="green")
    add_metric(slide, 10.15, 4.86, 2.2, "0 túnel", "backend en VPS", accent="coral")
    add_footer(slide, "Admin mobile", 6)
    return prs


def driver_deck():
    prs = deck_setup()
    # 1 Cover
    slide = base_slide(prs, dark=True)
    add_mark(slide, 0.7, 0.7, 0.66, dark=False)
    add_label(slide, "Presentación de producto", 1.55, 0.78, 4.0, color="amber_light")
    add_text(slide, "MiFlota Chofer", 0.7, 1.78, 7.3, 0.8, size=40, color="white", bold=True, margin=0)
    add_text(slide, "La operación diaria del conductor, sin vueltas.", 0.74, 2.75, 6.8, 0.54, size=18, color="amber_light", margin=0)
    add_pill(slide, "App móvil para conductores", 0.74, 3.65, 2.55, color="green_light", text_color="green", h=0.4, size=10)
    add_phone(slide, 8.35, 1.03, w=3.0, h=5.95, accent="green", driver=True)
    add_text(slide, "SALDO · PAGO · SOPORTE", 0.74, 6.52, 4.5, 0.2, size=9, color="green_light", bold=True, margin=0)
    add_footer(slide, "Chofer", 1)

    # 2 Value
    slide = base_slide(prs)
    add_title(slide, "01 · La oportunidad", "Más claridad para el chofer, menos idas y vueltas", "El conductor ve lo que debe, paga y deja constancia desde el celular.", 2)
    add_feature_card(slide, 0.72, 2.18, 3.72, 2.0, "01", "Saber el saldo", "Deuda pendiente, cuotas y próximo vencimiento en el inicio.", accent="coral")
    add_feature_card(slide, 4.8, 2.18, 3.72, 2.0, "02", "Pagar con evidencia", "Carga el importe y el comprobante para que el dueño lo valide.", accent="green")
    add_feature_card(slide, 8.88, 2.18, 3.72, 2.0, "03", "Pedir soporte", "Reporta una falla del vehículo con contexto y seguimiento.", accent="amber")
    add_phone(slide, 1.0, 4.8, w=2.15, h=2.15, accent="green", driver=True)
    add_text(slide, "Una relación más transparente empieza por la misma información.", 3.62, 5.38, 7.9, 0.45, size=20, color="ink", bold=True, margin=0)
    add_footer(slide, "Chofer", 2)

    # 3 Workflow
    slide = base_slide(prs, dark=True)
    add_title(slide, "02 · Flujo diario", "Del resumen al comprobante en menos de un minuto", "La app acompaña los momentos que más importan en la operación.", 3, dark=True)
    steps = [
        (0.82, "Inicio", "ver saldo", "coral"),
        (3.43, "Pagar", "importe + recibo", "green"),
        (6.04, "Pagos", "historial", "amber"),
        (8.65, "Reportar", "falla del auto", "coral"),
    ]
    for i, (x, title, body, accent) in enumerate(steps):
        add_box(slide, x, 3.05, 2.0, 1.38, "white", radius=True)
        add_pill(slide, f"0{i+1}", x + 0.2, 3.25, 0.44, color=accent, text_color="ink", h=0.29, size=8)
        add_text(slide, title, x + 0.2, 3.69, 1.55, 0.22, size=13, color="ink", bold=True, margin=0)
        add_text(slide, body, x + 0.2, 4.02, 1.55, 0.2, size=8.3, color="muted", margin=0)
        if i < len(steps) - 1:
            add_arrow(slide, x + 2.08, 3.74, x + 2.5, 3.74, color="amber")
    add_text(slide, "El chofer gana autonomía; el dueño gana trazabilidad.", 0.84, 5.38, 7.5, 0.3, size=17, color="amber_light", margin=0)
    add_pill(slide, "simple", 9.4, 5.23, 0.95, color="green_light", text_color="green", h=0.36, size=9)
    add_pill(slide, "visible", 10.55, 5.23, 0.95, color="amber_light", text_color="gold", h=0.36, size=9)
    add_pill(slide, "auditable", 11.7, 5.23, 1.05, color="coral_light", text_color="coral", h=0.36, size=9)
    add_footer(slide, "Chofer", 3)

    # 4 Screens
    slide = base_slide(prs)
    add_title(slide, "03 · Producto", "Todo lo necesario, sin ruido", "La navegación está centrada en saldo, pagos y ayuda.", 4)
    add_phone(slide, 0.78, 1.82, w=2.43, h=5.15, accent="coral", driver=True)
    add_phone(slide, 3.72, 1.82, w=2.43, h=5.15, accent="green", driver=True)
    add_phone(slide, 6.66, 1.82, w=2.43, h=5.15, accent="amber", driver=True)
    add_box(slide, 9.73, 2.05, 2.58, 4.7, "ink", radius=True)
    add_label(slide, "Momentos clave", 10.05, 2.4, 1.8, color="amber")
    add_rich_lines(slide, ["Inicio: resumen de cuenta.", "Pagar: comprobante.", "Pagos: historial claro.", "Reportes: falla del vehículo.", "Perfil: datos y acceso."], 10.03, 2.93, 1.95, 2.5, size=9.5, color="white", bullet_color="green_light", gap=0.48)
    add_footer(slide, "Chofer", 4)

    # 5 Trust / auth
    slide = base_slide(prs, dark=True)
    add_title(slide, "04 · Confianza", "Cada chofer ve sólo lo que necesita", "Acceso por sesión y alcance por vehículo para mantener la operación ordenada.", 5, dark=True)
    add_box(slide, 0.78, 2.25, 3.28, 2.9, "white", radius=True)
    add_pill(slide, "01", 1.12, 2.62, 0.45, color="green", text_color="ink", h=0.3, size=8)
    add_text(slide, "Inicio de sesión", 1.12, 3.08, 2.4, 0.25, size=14, color="ink", bold=True, margin=0)
    add_text(slide, "Usuario por nombre.apellido y contraseña entregada por el dueño.", 1.12, 3.48, 2.45, 0.7, size=10, color="muted", margin=0)
    add_box(slide, 4.65, 2.25, 3.28, 2.9, "white", radius=True)
    add_pill(slide, "02", 4.99, 2.62, 0.45, color="amber", text_color="ink", h=0.3, size=8)
    add_text(slide, "Cuenta acotada", 4.99, 3.08, 2.4, 0.25, size=14, color="ink", bold=True, margin=0)
    add_text(slide, "Saldo, pagos y reportes del vehículo asignado; sin acceso a la flota completa.", 4.99, 3.48, 2.5, 0.7, size=10, color="muted", margin=0)
    add_box(slide, 8.52, 2.25, 3.28, 2.9, "white", radius=True)
    add_pill(slide, "03", 8.86, 2.62, 0.45, color="coral", text_color="ink", h=0.3, size=8)
    add_text(slide, "Registro compartido", 8.86, 3.08, 2.55, 0.25, size=14, color="ink", bold=True, margin=0)
    add_text(slide, "El comprobante queda asociado al movimiento para que ambas partes trabajen con el mismo dato.", 8.86, 3.48, 2.5, 0.7, size=10, color="muted", margin=0)
    add_text(slide, "La confianza se diseña en el flujo, no se agrega al final.", 0.82, 5.75, 8.8, 0.3, size=16, color="amber_light", bold=True, margin=0)
    add_footer(slide, "Chofer", 5)

    # 6 Close
    slide = base_slide(prs)
    add_title(slide, "05 · Demo", "Un día más simple para cada conductor", "MiFlota Chofer convierte la administración en una conversación clara.", 6)
    add_phone(slide, 1.0, 1.82, w=2.6, h=5.1, accent="green", driver=True)
    add_box(slide, 4.55, 2.18, 7.55, 3.35, "ink", radius=True)
    add_label(slide, "Recorrido sugerido", 4.95, 2.58, 2.5, color="amber")
    add_rich_lines(slide, ["Ingresar con las credenciales entregadas.", "Revisar deuda y próximo vencimiento.", "Registrar un pago con comprobante.", "Abrir Pagos para verificar el historial.", "Reportar una falla si aparece."], 4.95, 3.1, 6.55, 1.75, size=12, color="white", bullet_color="green_light", gap=0.44)
    add_pill(slide, "claro para el chofer", 4.95, 5.05, 1.7, color="green_light", text_color="green", h=0.35, size=9)
    add_pill(slide, "útil para el dueño", 6.82, 5.05, 1.65, color="amber_light", text_color="gold", h=0.35, size=9)
    add_footer(slide, "Chofer", 6)
    return prs


def save(prs, filename):
    OUT.mkdir(parents=True, exist_ok=True)
    path = OUT / filename
    prs.save(path)
    return path


def main():
    paths = [
        save(web_deck(), "miflota-web.pptx"),
        save(mobile_deck(), "miflota-admin-mobile.pptx"),
        save(driver_deck(), "miflota-chofer.pptx"),
    ]
    for path in paths:
        prs = Presentation(path)
        print(f"{path}: {len(prs.slides)} slides, {path.stat().st_size:,} bytes")


if __name__ == "__main__":
    main()
